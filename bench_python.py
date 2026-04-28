#!/usr/bin/env python3
"""Benchmark Python document-parsing libraries against the test corpus.

Usage: .venv/bin/python bench_python.py ~/projects/office_oxide_tests/
"""

import json
import os
import resource
import sys
import time
import signal
from pathlib import Path

sys.setrecursionlimit(5000)

TIMEOUT = 30  # seconds per file


def peak_rss_kb() -> int:
    """Return peak RSS of the current process in KB.

    Linux reports ru_maxrss in KB; macOS in bytes. We only run bench on
    Linux, but normalize anyway so the harness doesn't silently lie.
    """
    raw = resource.getrusage(resource.RUSAGE_SELF).ru_maxrss
    if sys.platform == "darwin":
        return raw // 1024
    return raw

class TimeoutError(Exception):
    pass

def timeout_handler(signum, frame):
    raise TimeoutError("Timed out")

signal.signal(signal.SIGALRM, timeout_handler)

# ── Library wrappers ──────────────────────────────────────────────────────

def try_markitdown(path: str) -> str:
    from markitdown import MarkItDown
    md = MarkItDown()
    result = md.convert(path)
    return result.text_content

def try_python_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)

def try_openpyxl(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()
    return "\n".join(parts)

def try_python_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)

def try_python_calamine(path: str) -> str:
    from python_calamine import CalamineWorkbook
    wb = CalamineWorkbook.from_path(path)
    parts = []
    for name in wb.sheet_names:
        sheet = wb.get_sheet_by_name(name)
        for row in sheet.to_python():
            parts.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(parts)

def try_xlrd(path: str) -> str:
    import xlrd
    book = xlrd.open_workbook(path, on_demand=True)
    parts = []
    for sheet in book.sheets():
        for rx in range(sheet.nrows):
            row = sheet.row_values(rx)
            parts.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(parts)


# ── Runner ────────────────────────────────────────────────────────────────

LIBS = {
    ".docx": [
        ("markitdown", try_markitdown),
        ("python-docx", try_python_docx),
    ],
    ".xlsx": [
        ("markitdown", try_markitdown),
        ("openpyxl", try_openpyxl),
        ("python-calamine", try_python_calamine),
    ],
    ".pptx": [
        ("markitdown", try_markitdown),
        ("python-pptx", try_python_pptx),
    ],
    ".xls": [
        ("python-calamine", try_python_calamine),
        ("xlrd", try_xlrd),
    ],
}


def collect_files(root: str) -> list[str]:
    files = []
    for dirpath, _, filenames in os.walk(root):
        for f in filenames:
            ext = os.path.splitext(f)[1].lower()
            if ext in (".docx", ".xlsx", ".pptx", ".xls"):
                files.append(os.path.join(dirpath, f))
    files.sort()
    return files


def main():
    if len(sys.argv) < 2:
        print("Usage: bench_python.py DIR [--json OUT]", file=sys.stderr)
        sys.exit(1)

    root = sys.argv[1]
    json_out: str | None = None
    if "--json" in sys.argv:
        i = sys.argv.index("--json")
        json_out = sys.argv[i + 1]

    files = collect_files(root)
    print(f"Found {len(files)} files", file=sys.stderr)
    rss_start_kb = peak_rss_kb()

    # stats[lib_name] = {"ok": 0, "fail": 0, "errors": {}, "total_ms": 0}
    stats: dict[str, dict] = {}

    for i, path in enumerate(files):
        ext = os.path.splitext(path)[1].lower()
        libs_for_ext = LIBS.get(ext, [])

        for lib_name, func in libs_for_ext:
            if lib_name not in stats:
                stats[lib_name] = {"ok": 0, "fail": 0, "errors": {}, "total_ms": 0.0, "timeout": 0}

            signal.alarm(TIMEOUT)
            t0 = time.monotonic()
            try:
                text = func(path)
                elapsed = (time.monotonic() - t0) * 1000
                stats[lib_name]["ok"] += 1
                stats[lib_name]["total_ms"] += elapsed
            except TimeoutError:
                stats[lib_name]["fail"] += 1
                stats[lib_name]["timeout"] += 1
                err_key = "TIMEOUT"
                stats[lib_name]["errors"][err_key] = stats[lib_name]["errors"].get(err_key, 0) + 1
            except BaseException as e:
                elapsed = (time.monotonic() - t0) * 1000
                stats[lib_name]["fail"] += 1
                stats[lib_name]["total_ms"] += elapsed
                err_key = type(e).__name__
                stats[lib_name]["errors"][err_key] = stats[lib_name]["errors"].get(err_key, 0) + 1
            finally:
                signal.alarm(0)

        if (i + 1) % 500 == 0:
            print(f"  [{i+1}/{len(files)}]", file=sys.stderr)

    rss_end_kb = peak_rss_kb()

    # Print results
    print("\n=== Python Library Benchmark Results ===\n")
    print(f"Peak RSS (harness process): {rss_end_kb / 1024:.1f} MiB "
          f"(delta {(rss_end_kb - rss_start_kb) / 1024:.1f} MiB)\n")
    for lib_name, s in sorted(stats.items()):
        total = s["ok"] + s["fail"]
        pct = s["ok"] / total * 100 if total > 0 else 0
        wall = s["total_ms"] / 1000
        print(f"{lib_name}:")
        print(f"  Total: {total}  OK: {s['ok']}  FAIL: {s['fail']}  Rate: {pct:.1f}%  Wall: {wall:.1f}s")
        if s["errors"]:
            top_errors = sorted(s["errors"].items(), key=lambda x: -x[1])[:10]
            for err, count in top_errors:
                print(f"    {err}: {count}")
        print()

    if json_out:
        payload = {
            "peak_rss_kb": rss_end_kb,
            "peak_rss_delta_kb": rss_end_kb - rss_start_kb,
            "libs": stats,
        }
        with open(json_out, "w") as f:
            json.dump(payload, f, indent=2, sort_keys=True)
        print(f"Wrote JSON results to {json_out}", file=sys.stderr)


if __name__ == "__main__":
    main()
